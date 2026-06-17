from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY    = RGBColor(0x1A, 0x60, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG   = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)

prs = Presentation("MediScript-E-Presentation.pptx")

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC",
       0.3, 7.22, 12.7, 0.25, size=9, color=WHITE, align=PP_ALIGN.CENTER)

# ── Rebuild Slide 3 — Outline (Abstract removed, renumbered) ──
slide3 = prs.slides[2]
for shape in list(slide3.shapes):
    shape._element.getparent().remove(shape._element)

# Header
rect(slide3, 0, 0, 13.33, 1.25, fill=PRIMARY)
tb(slide3, "Presentation Outline", 0.4, 0.12, 12, 0.7,
   size=26, bold=True, color=WHITE)
tb(slide3, "What We'll Cover Today",
   0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

# Outline items — Abstract removed, renumbered 1–13
outline_left = [
    ("1",  "Related Work",           "IEEE"),
    ("2",  "Problem Statement",      "SDLC Phase 1"),
    ("3",  "Proposed Solution",      "SDLC Phase 1"),
    ("4",  "Requirements Analysis",  "SDLC Phase 1"),
    ("5",  "Use Cases & Actors",     "SDLC Phase 1"),
    ("6",  "System Architecture",    "SDLC Phase 2"),
    ("7",  "Database Design",        "SDLC Phase 2"),
]
outline_right = [
    ("8",  "Tech Stack",             "SDLC Phase 3"),
    ("9",  "Auth & Security",        "SDLC Phase 3"),
    ("10", "Core Features",          "SDLC Phase 3"),
    ("11", "Testing",                "SDLC Phase 4"),
    ("12", "Results & Deployment",   "SDLC Phase 5"),
    ("13", "Challenges & Conclusion","IEEE"),
]

for i, (num, title, phase) in enumerate(outline_left):
    y = 1.45 + i * 0.82
    rect(slide3, 0.4, y, 0.5, 0.62, fill=PRIMARY)
    tb(slide3, num, 0.4, y+0.12, 0.5, 0.38,
       size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide3, 0.95, y, 5.5, 0.62,
         fill=LIGHT_BG if i % 2 == 0 else WHITE)
    tb(slide3, title, 1.05, y+0.12, 3.5, 0.38, size=13, color=DARK)
    tb(slide3, phase, 4.55, y+0.12, 1.8, 0.38,
       size=10, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT)

for i, (num, title, phase) in enumerate(outline_right):
    y = 1.45 + i * 0.82
    rect(slide3, 6.9, y, 0.5, 0.62, fill=PRIMARY)
    tb(slide3, num, 6.9, y+0.12, 0.5, 0.38,
       size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide3, 7.45, y, 5.5, 0.62,
         fill=LIGHT_BG if i % 2 == 0 else WHITE)
    tb(slide3, title, 7.55, y+0.12, 3.5, 0.38, size=13, color=DARK)
    tb(slide3, phase, 11.05, y+0.12, 1.8, 0.38,
       size=10, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT)

footer(slide3)

prs.save("MediScript-E-Presentation.pptx")
print("✅ Slide 3 rebuilt — Abstract removed, renumbered 1–13")

# Verify
prs2 = Presentation("MediScript-E-Presentation.pptx")
slide3_check = prs2.slides[2]
texts = [s.text_frame.text.strip() for s in slide3_check.shapes
         if s.has_text_frame and s.text_frame.text.strip()]
print("\nSlide 3 content:")
for t in texts:
    if len(t) > 1:
        print(f"  {t}")
