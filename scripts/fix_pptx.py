from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY    = RGBColor(0x1A, 0x60, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
GRAY       = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG   = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT     = RGBColor(0x0D, 0x4A, 0x63)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)

prs = Presentation("MediScript-E-Presentation.pptx")
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def lines(slide, items, x, y, w, h, size=15, color=DARK,
          bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return t

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.25, fill=PRIMARY)
    tb(slide, title, 0.4, 0.12, 12, 0.7, size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC",
       0.3, 7.22, 12.7, 0.25, size=9, color=WHITE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    tbl = slide.shapes.add_table(
        len(rows)+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)
    ).table
    cw = Inches(w / cols)
    for i in range(cols): tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.color.rgb = DARK
    return tbl

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide = slides[old_index]
    xml_slides.remove(slide)
    xml_slides.insert(new_index, slide)

# ── FIX 1: Remove orphan placeholder from Slide 4 (Problem Statement) ──
slide4 = prs.slides[3]
to_remove = []
for shape in slide4.shapes:
    if shape.has_text_frame:
        full_text = shape.text_frame.text
        if "📸" in full_text or "ADD IMAGE" in full_text or "illustration" in full_text:
            to_remove.append(shape)
for shape in to_remove:
    shape._element.getparent().remove(shape._element)
print(f"✅ Fix 1: Removed {len(to_remove)} orphan placeholder(s) from Slide 4")

# ── FIX 2: Rebuild Slide 3 (Outline) with correct SDLC phase mapping ──
slide3 = prs.slides[2]
# Remove all existing shapes
for shape in list(slide3.shapes):
    shape._element.getparent().remove(shape._element)

# Rebuild cleanly
header(slide3, "Presentation Outline", "What We'll Cover Today")

outline_left = [
    ("1",  "Abstract",               "IEEE"),
    ("2",  "Related Work",           "IEEE"),
    ("3",  "Problem Statement",      "SDLC Phase 1"),
    ("4",  "Proposed Solution",      "SDLC Phase 1"),
    ("5",  "Requirements Analysis",  "SDLC Phase 1"),
    ("6",  "Use Cases & Actors",     "SDLC Phase 1"),
    ("7",  "System Architecture",    "SDLC Phase 2"),
]
outline_right = [
    ("8",  "Database Design",        "SDLC Phase 2"),
    ("9",  "Tech Stack",             "SDLC Phase 3"),
    ("10", "Auth & Security",        "SDLC Phase 3"),
    ("11", "Core Features",          "SDLC Phase 3"),
    ("12", "Testing",                "SDLC Phase 4"),
    ("13", "Results & Deployment",   "SDLC Phase 5"),
    ("14", "Challenges & Conclusion","IEEE"),
]

for i, (num, title, phase) in enumerate(outline_left):
    y = 1.45 + i * 0.77
    rect(slide3, 0.4, y, 0.5, 0.58, fill=PRIMARY)
    tb(slide3, num, 0.4, y+0.1, 0.5, 0.38,
       size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide3, 0.95, y, 5.5, 0.58,
         fill=LIGHT_BG if i % 2 == 0 else WHITE)
    tb(slide3, title, 1.05, y+0.1, 3.5, 0.38, size=13, color=DARK)
    tb(slide3, phase, 4.55, y+0.1, 1.8, 0.38,
       size=10, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT)

for i, (num, title, phase) in enumerate(outline_right):
    y = 1.45 + i * 0.77
    rect(slide3, 6.9, y, 0.5, 0.58, fill=PRIMARY)
    tb(slide3, num, 6.9, y+0.1, 0.5, 0.38,
       size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide3, 7.45, y, 5.5, 0.58,
         fill=LIGHT_BG if i % 2 == 0 else WHITE)
    tb(slide3, title, 7.55, y+0.1, 3.5, 0.38, size=13, color=DARK)
    tb(slide3, phase, 11.05, y+0.1, 1.8, 0.38,
       size=10, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT)

footer(slide3)
print("✅ Fix 2: Slide 3 (Outline) rebuilt with correct SDLC + IEEE phase labels")

# ── FIX 3: Add Related Work slide (IEEE requirement) ──
rw = prs.slides.add_slide(BLANK)
header(rw, "Related Work", "IEEE — Existing Systems & Gap Analysis")

table(rw,
    ["Existing System", "What It Does", "What's Missing"],
    [
        ["Practo (India)",
         "Doctor discovery & appointment booking",
         "No prescription management, no medicine reminders, no medical vault"],
        ["Zocdoc (USA)",
         "Online appointment scheduling",
         "No e-prescriptions, no 2FA, no AI chatbot, no file storage"],
        ["HealthTap",
         "AI symptom checker + doctor chat",
         "No appointment booking, no prescription PDF, no RBAC"],
        ["Manual / Paper-based",
         "Traditional in-person healthcare",
         "No digitization — prescriptions lost, records scattered, no reminders"],
    ],
    0.4, 1.45, 12.5, 3.6
)

rect(rw, 0.4, 5.2, 12.5, 1.7, fill=LIGHT_BG)
tb(rw, "Gap Identified", 0.6, 5.3, 12.1, 0.38,
   size=15, bold=True, color=PRIMARY)
lines(rw, [
    "•  No existing free platform combines: appointment booking + e-prescription + medicine reminders + medical vault + AI chatbot",
    "•  Most platforms lack role-based access control (RBAC) and Two-Factor Authentication (2FA)",
    "•  MediScript-E fills this gap — open, secure, full-stack, and production-deployed",
], 0.6, 5.7, 12.1, 1.1, size=13, color=DARK)

footer(rw)

# Move Related Work to index 3 (after Abstract at index 1, Outline at index 2)
move_slide(prs, len(prs.slides)-1, 3)
print("✅ Fix 3: Related Work slide added at position 4 (after Outline)")

prs.save("MediScript-E-Presentation.pptx")
print(f"\n✅ Total slides: {len(prs.slides)}")
print("🎉 All 3 fixes applied — MediScript-E-Presentation.pptx updated!")
