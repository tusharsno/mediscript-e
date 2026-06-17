from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY    = RGBColor(0x1A, 0x60, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
GRAY       = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG   = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT     = RGBColor(0x0D, 0x4A, 0x63)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)

prs = Presentation("MediScript-E-Presentation.pptx")

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def lines(slide, items, x, y, w, h, size=15, color=DARK,
          bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return t

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC",
       0.3, 7.22, 12.7, 0.25, size=9, color=WHITE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    tbl = slide.shapes.add_table(
        len(rows)+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)
    ).table
    cw = Inches(w / cols)
    for i in range(cols): tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.color.rgb = DARK
    return tbl

# ── FIX 1: Slide 2 — Change subtitle from "Phase 0" to "IEEE" ──
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Phase 0" in run.text:
                    run.text = run.text.replace("Phase 0 — Project Overview", "IEEE — Project Overview")
print("✅ Fix 1: Slide 2 subtitle updated — Phase 0 → IEEE")

# ── FIX 2: Slide 4 — Rebuild Related Work with proper table ──
slide4 = prs.slides[3]
for shape in list(slide4.shapes):
    shape._element.getparent().remove(shape._element)

# Rebuild header
rect(slide4, 0, 0, 13.33, 1.25, fill=PRIMARY)
tb(slide4, "Related Work", 0.4, 0.12, 12, 0.7,
   size=26, bold=True, color=WHITE)
tb(slide4, "IEEE — Existing Systems & Gap Analysis",
   0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

# Table
table(slide4,
    ["Existing System", "What It Does", "What's Missing"],
    [
        ["Practo (India)",
         "Doctor discovery & appointment booking",
         "No prescription mgmt, no reminders, no vault"],
        ["Zocdoc (USA)",
         "Online appointment scheduling",
         "No e-prescriptions, no 2FA, no AI chatbot"],
        ["HealthTap",
         "AI symptom checker + doctor chat",
         "No appointment booking, no RBAC, no PDF"],
        ["Manual / Paper-based",
         "Traditional in-person healthcare",
         "No digitization — records lost, no reminders"],
    ],
    0.4, 1.45, 12.5, 3.5
)

# Gap section
rect(slide4, 0.4, 5.1, 12.5, 1.85, fill=LIGHT_BG)
tb(slide4, "Gap Identified", 0.6, 5.2, 12.1, 0.38,
   size=15, bold=True, color=PRIMARY)
lines(slide4, [
    "•  No free platform combines: booking + e-prescription + reminders + vault + AI chatbot + 2FA",
    "•  Most platforms lack Role-Based Access Control (RBAC) and Two-Factor Authentication",
    "•  MediScript-E fills this gap — open, secure, full-stack, production-deployed on Vercel",
], 0.6, 5.62, 12.1, 1.2, size=13, color=DARK)

footer(slide4)
print("✅ Fix 2: Slide 4 Related Work rebuilt with proper table content")

# ── FIX 3: Slide 19 — Rebuild Challenges with proper table ──
slide19 = prs.slides[18]
for shape in list(slide19.shapes):
    shape._element.getparent().remove(shape._element)

# Rebuild header
rect(slide19, 0, 0, 13.33, 1.25, fill=PRIMARY)
tb(slide19, "Challenges & Solutions", 0.4, 0.12, 12, 0.7,
   size=26, bold=True, color=WHITE)
tb(slide19, "Real Problems — Real Fixes",
   0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

table(slide19,
    ["#", "Challenge", "Solution"],
    [
        ["1", "2FA re-auth after OTP verification",
              "Magic token + DB OTP null check before session creation"],
        ["2", "Supabase SSL failure on Vercel",
              "sslmode=no-verify + PgBouncer pooler URL"],
        ["3", "OAuth users have no password field",
              "Check user.password === '' before allowing update"],
        ["4", "Medicine reminder time zone mismatch",
              "Convert server time to Bangladesh Standard Time (UTC+6)"],
        ["5", "Prisma connection exhaustion in serverless",
              "Singleton pattern via globalThis.prismaGlobal"],
        ["6", "Cascade delete — foreign key constraint errors",
              "Manual delete in correct dependency order"],
        ["7", "Stale JWT after profile name update",
              "Force window.location.reload() after save"],
    ],
    0.4, 1.45, 12.5, 5.6
)

footer(slide19)
print("✅ Fix 3: Slide 19 Challenges rebuilt with proper table content")

prs.save("MediScript-E-Presentation.pptx")
print("\n🎉 All 3 fixes applied and saved.")
