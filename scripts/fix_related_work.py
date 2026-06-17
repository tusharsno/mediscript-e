from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY    = RGBColor(0x1A, 0x60, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG   = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT     = RGBColor(0x0D, 0x4A, 0x63)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)

prs = Presentation("MediScript-E-Presentation.pptx")

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def lines(slide, items, x, y, w, h, size=15, color=DARK, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.color.rgb = color
    return t

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC",
       0.3, 7.22, 12.7, 0.25, size=9, color=WHITE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    tbl = slide.shapes.add_table(
        len(rows)+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)
    ).table
    cw = Inches(w / cols)
    for i in range(cols): tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.color.rgb = DARK
    return tbl

# ── Rebuild Slide 4 — Related Work ──
slide4 = prs.slides[3]
for shape in list(slide4.shapes):
    shape._element.getparent().remove(shape._element)

rect(slide4, 0, 0, 13.33, 1.25, fill=PRIMARY)
tb(slide4, "Related Work", 0.4, 0.12, 12, 0.7,
   size=26, bold=True, color=WHITE)
tb(slide4, "IEEE — Existing Systems & Gap Analysis",
   0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

table(slide4,
    ["Platform", "Primary Focus", "What It Lacks (Combined)"],
    [
        ["Practo (India)",
         "Doctor discovery & appointment booking",
         "No unified: reminders + vault + AI chatbot + open access"],
        ["Zocdoc (USA)",
         "Online appointment scheduling",
         "No unified: e-prescriptions + 2FA + AI + free open platform"],
        ["HealthTap (USA)",
         "AI symptom checker + doctor chat",
         "No unified: appointment booking + RBAC + prescription PDF"],
        ["Manual / Paper-based",
         "Traditional in-person healthcare",
         "No digitization at all — records lost, no reminders"],
    ],
    0.4, 1.45, 12.5, 3.4
)

rect(slide4, 0.4, 5.0, 12.5, 1.95, fill=LIGHT_BG)
tb(slide4, "The Gap", 0.6, 5.1, 12.1, 0.38,
   size=15, bold=True, color=PRIMARY)
lines(slide4, [
    "•  Each platform solves one problem — none combines all features in a single free, open system",
    "•  No existing free platform offers: booking + e-prescription + reminders + vault + AI chatbot + 2FA + RBAC",
    "•  MediScript-E fills this gap — unified, secure, role-based, and production-deployed",
], 0.6, 5.52, 12.1, 1.3, size=13, color=DARK)

footer(slide4)
prs.save("MediScript-E-Presentation.pptx")
print("✅ Slide 4 rebuilt with safe, accurate claims")
