from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY  = RGBColor(0x1A, 0x60, 0x80)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1E, 0x29, 0x3B)
GRAY     = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT   = RGBColor(0x0D, 0x4A, 0x63)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return t

def lines(slide, items, x, y, w, h, size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return t

def bullets(slide, items, x, y, w, h, size=15, color=DARK):
    return lines(slide, [f"•  {i}" for i in items], x, y, w, h, size=size, color=color)

def img_placeholder(slide, x, y, w, h, label="ADD IMAGE"):
    b = rect(slide, x, y, w, h, fill=LIGHT_BG, line=PRIMARY)
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"📸  {label}"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    return b

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.25, fill=PRIMARY)
    tb(slide, title, 0.4, 0.12, 12, 0.7, size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC", 0.3, 7.22, 12.7, 0.25,
       size=9, color=WHITE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    cw = Inches(w / cols)
    for i in range(cols): tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.color.rgb = DARK
    return tbl

# ── SLIDE 1: Title ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=PRIMARY)
rect(s, 0, 0, 7.0, 7.5, fill=ACCENT)
tb(s, "MediScript-E", 0.5, 1.4, 6.2, 1.1, size=50, bold=True, color=WHITE)
tb(s, "A Secure Digital Healthcare Platform", 0.5, 2.6, 6.2, 0.55, size=18, color=LIGHT_BLUE)
tb(s, "Software Engineering Course Project", 0.5, 3.2, 6.2, 0.45, size=13, italic=True, color=LIGHT_BLUE)
rect(s, 0.5, 3.85, 2.5, 0.05, fill=WHITE)
lines(s, ["Tushar Barua", "CSE — USTC, Chittagong", "2024–2025"], 0.5, 4.05, 5.5, 1.1, size=14, color=WHITE)
tb(s, "🌐  mediscript-e.vercel.app", 0.5, 5.35, 5.5, 0.4, size=12, color=LIGHT_BLUE)
img_placeholder(s, 7.3, 1.0, 5.7, 5.5, "Landing page screenshot")
print("✅ Slide 1")

# ── SLIDE 2: Abstract ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Abstract", "Phase 0 — Project Overview")
rect(s, 0.4, 1.45, 8.0, 4.5, fill=LIGHT_BG)
tb(s, "What is MediScript-E?", 0.6, 1.6, 7.5, 0.45, size=17, bold=True, color=PRIMARY)
bullets(s, [
    "Full-stack digital healthcare platform",
    "Connects Patients, Doctors, and Administrators",
    "Built following complete SDLC methodology",
    "Digitizes the entire patient-doctor workflow",
], 0.6, 2.1, 7.5, 1.8, size=15)
tb(s, "Core Purpose", 0.6, 3.95, 7.5, 0.4, size=15, bold=True, color=PRIMARY)
bullets(s, [
    "Digitize patient-doctor workflow end-to-end",
    "Secure medical data in the cloud",
    "Automate medicine reminders via email",
], 0.6, 4.38, 7.5, 1.1, size=14)
tb(s, "Built with:  Next.js  ·  TypeScript  ·  PostgreSQL  ·  Vercel", 0.6, 5.55, 7.5, 0.4, size=13, italic=True, color=GRAY)
img_placeholder(s, 8.8, 1.45, 4.1, 4.5, "Dashboard overview screenshot")
footer(s)
print("✅ Slide 2")

# ── SLIDE 3: Problem Statement ──────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Problem Statement", "Phase 1 — Requirements Analysis")
table(s,
    ["Pain Point", "Reality"],
    [
        ["Prescriptions",    "Lost, damaged, inaccessible"],
        ["Appointments",     "Phone calls & physical visits required"],
        ["Medical Records",  "Scattered across multiple hospitals"],
        ["Medicine Schedule","No reminders — often forgotten"],
        ["Data Security",    "No role-based access control"],
    ],
    0.5, 1.45, 8.5, 3.5
)
rect(s, 0.5, 5.1, 8.5, 0.9, fill=ACCENT)
tb(s, '"Healthcare data is fragmented, insecure, and inaccessible when needed most."',
   0.7, 5.18, 8.1, 0.75, size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
img_placeholder(s, 9.3, 1.45, 3.7, 4.55, "Problem illustration / icon")
footer(s)
print("✅ Slide 3")

# ── SLIDE 4: Proposed Solution ──────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Proposed Solution", "Phase 1 — Requirements Analysis")
tb(s, "One Platform. Three Roles. Complete Workflow.", 0.5, 1.4, 12.5, 0.5, size=18, bold=True, color=PRIMARY)
table(s,
    ["Role", "Key Capabilities"],
    [
        ["👤  Patient", "Book appointments · View prescriptions · Set reminders · Upload records"],
        ["👨‍⚕️  Doctor",  "Manage appointments · Issue digital prescriptions · Archive records"],
        ["🛡️  Admin",   "Monitor platform · Manage users · View all appointments"],
    ],
    0.5, 2.0, 12.3, 2.5
)
tb(s, "Result:", 0.5, 4.7, 1.2, 0.4, size=15, bold=True, color=PRIMARY)
tb(s, "End-to-end digital healthcare — from booking to prescription to secure storage.",
   1.7, 4.7, 11.0, 0.4, size=15, color=DARK)
rect(s, 0.5, 5.2, 12.3, 1.5, fill=LIGHT_BG)
cols = [
    ("25+", "API Endpoints"),
    ("9", "DB Models"),
    ("3", "User Roles"),
    ("5", "Core Features"),
    ("100%", "Live on Vercel"),
]
for i, (num, label) in enumerate(cols):
    x = 0.7 + i * 2.45
    tb(s, num,   x, 5.3,  2.2, 0.6, size=28, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    tb(s, label, x, 5.88, 2.2, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)
footer(s)
print("✅ Slide 4")

# ── SLIDE 5: Requirements ───────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Requirements Analysis", "Phase 1 — Functional & Non-Functional")
# Left column
rect(s, 0.4, 1.4, 6.0, 5.5, fill=LIGHT_BG)
tb(s, "Functional Requirements", 0.6, 1.5, 5.6, 0.45, size=15, bold=True, color=PRIMARY)
bullets(s, [
    "FR1: Register & login with email verification",
    "FR2: Book appointments with available doctors",
    "FR3: View & download prescriptions as PDF",
    "FR4: Set medicine reminders — automated email alerts",
    "FR5: Upload & manage medical records (cloud)",
    "FR6: Doctor manages appointments & prescriptions",
    "FR7: Admin monitors users & platform stats",
], 0.6, 2.0, 5.6, 3.5, size=13)
# Right column
rect(s, 6.9, 1.4, 6.0, 5.5, fill=LIGHT_BG)
tb(s, "Non-Functional Requirements", 7.1, 1.5, 5.6, 0.45, size=15, bold=True, color=PRIMARY)
table(s,
    ["NFR", "Requirement"],
    [
        ["Security",       "2FA, RBAC, bcrypt, JWT"],
        ["Performance",    "Serverless, connection pooling"],
        ["Scalability",    "Vercel + Supabase managed infra"],
        ["Usability",      "Responsive, mobile-friendly UI"],
        ["Reliability",    "Token expiry, cascade delete"],
        ["Maintainability","TypeScript, Prisma ORM"],
    ],
    7.1, 2.0, 5.6, 4.6
)
footer(s)
print("✅ Slide 5")

prs.save("MediScript-E-Presentation.pptx")
print("✅ Slides 1–5 saved.")
