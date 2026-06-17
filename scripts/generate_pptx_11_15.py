from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY  = RGBColor(0x1A, 0x60, 0x80)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1E, 0x29, 0x3B)
GRAY     = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT   = RGBColor(0x0D, 0x4A, 0x63)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation("MediScript-E-Presentation.pptx")
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def lines(slide, items, x, y, w, h, size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return t

def bullets(slide, items, x, y, w, h, size=15, color=DARK):
    return lines(slide, [f"•  {i}" for i in items], x, y, w, h, size=size, color=color)

def img_placeholder(slide, x, y, w, h, label="ADD IMAGE"):
    b = rect(slide, x, y, w, h, fill=LIGHT_BG, line=PRIMARY)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"📸  {label}"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = PRIMARY
    return b

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.25, fill=PRIMARY)
    tb(slide, title, 0.4, 0.12, 12, 0.7, size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC", 0.3, 7.22, 12.7, 0.25,
       size=9, color=WHITE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    cw = Inches(w / cols)
    for i in range(cols): tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.color.rgb = DARK
    return tbl

# ── SLIDE 11: Patient Features ───────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Patient Features", "Phase 3 — Implementation")
features = [
    ("📅", "Appointment Booking",  "Select doctor → pick date & time → confirm booking"),
    ("📄", "E-Prescription",       "View diagnosis + medications → download as PDF"),
    ("💊", "Medicine Reminders",   "Set schedule → receive automated email alerts daily"),
    ("🗂️", "Medical Vault",        "Upload reports → stored securely in Supabase cloud"),
]
for i, (icon, title, desc) in enumerate(features):
    x = 0.4 + (i % 2) * 6.45
    y = 1.4 + (i // 2) * 2.65
    rect(s, x, y, 6.1, 2.4, fill=LIGHT_BG)
    tb(s, icon, x+0.2, y+0.15, 0.7, 0.7, size=28, align=PP_ALIGN.CENTER)
    tb(s, title, x+1.0, y+0.2, 4.8, 0.45, size=16, bold=True, color=PRIMARY)
    tb(s, desc,  x+1.0, y+0.7, 4.8, 0.5,  size=13, color=GRAY)
    img_placeholder(s, x+0.2, y+1.3, 5.7, 0.85, f"{title} screenshot")
footer(s)
print("✅ Slide 11")

# ── SLIDE 12: Doctor & Admin Features ───────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Doctor & Admin Features", "Phase 3 — Implementation")
rect(s, 0.4, 1.4, 6.1, 5.5, fill=LIGHT_BG)
tb(s, "👨⚕️  Doctor", 0.6, 1.5, 5.7, 0.45, size=17, bold=True, color=PRIMARY)
bullets(s, [
    "Confirm / cancel / complete appointments",
    "Issue prescriptions → appointment auto-completed",
    "Edit, archive, delete prescriptions",
    "See patient blood group on prescription form",
], 0.6, 2.05, 5.7, 2.0, size=14)
img_placeholder(s, 0.6, 4.15, 5.7, 2.5, "Doctor dashboard screenshot")

rect(s, 6.9, 1.4, 6.1, 5.5, fill=LIGHT_BG)
tb(s, "🛡️  Admin", 7.1, 1.5, 5.7, 0.45, size=17, bold=True, color=PRIMARY)
bullets(s, [
    "Real-time platform statistics",
    "Delete users — cascade removes all data",
    "Monitor all appointments & contact messages",
], 7.1, 2.05, 5.7, 1.5, size=14)
img_placeholder(s, 7.1, 3.65, 5.7, 3.0, "Admin dashboard screenshot")
footer(s)
print("✅ Slide 12")

# ── SLIDE 13: MediBot & Search ───────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "MediBot & Global Search", "Phase 3 — Special Features")
rect(s, 0.4, 1.4, 6.1, 5.5, fill=LIGHT_BG)
tb(s, "🤖  MediBot — AI Chatbot", 0.6, 1.5, 5.7, 0.45, size=16, bold=True, color=PRIMARY)
bullets(s, [
    "Powered by Groq — Llama 3.1 8B Instant",
    "Answers platform-related questions",
    "Available globally on all pages",
    "Responds in Bangla if user writes in Bangla",
    "Max 512 tokens — concise responses",
], 0.6, 2.05, 5.7, 2.3, size=13)
img_placeholder(s, 0.6, 4.45, 5.7, 2.2, "MediBot chatbot screenshot")

rect(s, 6.9, 1.4, 6.1, 5.5, fill=LIGHT_BG)
tb(s, "🔍  Global Search", 7.1, 1.5, 5.7, 0.45, size=16, bold=True, color=PRIMARY)
bullets(s, [
    "Trigger: navbar icon or Ctrl+K",
    "Debounced — 300ms delay",
], 7.1, 2.05, 5.7, 0.9, size=13)
table(s,
    ["Role", "Can Search"],
    [
        ["Guest",   "Doctors only"],
        ["Patient", "Doctors · own appointments · prescriptions"],
        ["Doctor",  "Patients · own appointments · prescriptions"],
        ["Admin",   "All users · appointments · contacts"],
    ],
    7.1, 3.1, 5.7, 2.8
)
footer(s)
print("✅ Slide 13")

# ── SLIDE 14: Testing ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Testing", "Phase 4 — Testing & Validation")
tb(s, "Testing Approach:  Manual + Integration Testing", 0.5, 1.4, 12.3, 0.4, size=15, bold=True, color=PRIMARY)
table(s,
    ["Test Area", "Coverage", "Result"],
    [
        ["Auth flows (register, login, 2FA, OAuth)", "All auth paths tested end-to-end",       "✅ Pass"],
        ["Appointment lifecycle",                    "PENDING → CONFIRMED → COMPLETED",         "✅ Pass"],
        ["Prescription CRUD",                        "Create, edit, archive, delete",           "✅ Pass"],
        ["Medicine reminder + email delivery",       "Cron trigger → email received",           "✅ Pass"],
        ["Medical vault upload/delete",              "Upload, rename, delete verified",         "✅ Pass"],
        ["Security & authorization",                 "Unauthorized access → 401/403 verified",  "✅ Pass"],
        ["Edge cases",                               "Past date, duplicate email, expired OTP", "✅ Pass"],
    ],
    0.5, 1.95, 12.3, 5.0
)
footer(s)
print("✅ Slide 14")

# ── SLIDE 15: Deployment ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Deployment", "Phase 5 — Deployment")
steps = [
    "Push to GitHub (main branch)",
    "Vercel detects push → triggers build",
    "Next.js build (pnpm build)",
    "Serverless functions deployed globally",
    "Live at: mediscript-e.vercel.app  ✅",
]
for i, step in enumerate(steps):
    y = 1.45 + i * 0.72
    rect(s, 0.4, y, 7.5, 0.58, fill=LIGHT_BG if i % 2 == 0 else WHITE)
    rect(s, 0.4, y, 0.55, 0.58, fill=PRIMARY)
    tb(s, str(i+1), 0.4, y+0.1, 0.55, 0.38, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, 1.05, y+0.12, 6.7, 0.35, size=13, color=DARK)
    if i < 4:
        tb(s, "↓", 0.55, y+0.58, 0.55, 0.3, size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

tb(s, "Key Config", 0.4, 5.15, 7.5, 0.38, size=14, bold=True, color=PRIMARY)
bullets(s, [
    "Database: Supabase PostgreSQL + PgBouncer pooling",
    "Cron: GitHub Actions — every 5 min → reminder emails",
    "10 environment variables configured on Vercel",
], 0.4, 5.55, 7.5, 1.3, size=13)
img_placeholder(s, 8.3, 1.45, 4.7, 5.5, "Vercel deployment dashboard screenshot")
footer(s)
print("✅ Slide 15")

prs.save("MediScript-E-Presentation.pptx")
print("✅ Slides 11–15 saved.")
