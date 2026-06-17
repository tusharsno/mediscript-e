from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY  = RGBColor(0x1A, 0x60, 0x80)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1E, 0x29, 0x3B)
GRAY     = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT   = RGBColor(0x0D, 0x4A, 0x63)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation("MediScript-E-Presentation.pptx")
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def lines(slide, items, x, y, w, h, size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return t

def bullets(slide, items, x, y, w, h, size=15, color=DARK):
    return lines(slide, [f"•  {i}" for i in items], x, y, w, h, size=size, color=color)

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.25, fill=PRIMARY)
    tb(slide, title, 0.4, 0.12, 12, 0.7, size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC", 0.3, 7.22, 12.7, 0.25,
       size=9, color=WHITE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    cw = Inches(w / cols)
    for i in range(cols): tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.color.rgb = DARK
    return tbl

# ── SLIDE 16: Challenges & Solutions ────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Challenges & Solutions", "Real Problems — Real Fixes")
table(s,
    ["#", "Challenge", "Solution"],
    [
        ["1", "2FA re-auth after OTP",              "Magic token + DB OTP null check before session"],
        ["2", "Supabase SSL failure on Vercel",     "sslmode=no-verify + PgBouncer pooler URL"],
        ["3", "OAuth users — no password field",    "Check empty password before allowing update"],
        ["4", "Reminder time zone mismatch",        "Convert server time to UTC+6 (Bangladesh)"],
        ["5", "Prisma connection in serverless",    "Singleton pattern via globalThis"],
        ["6", "Cascade delete constraint errors",   "Manual delete in correct dependency order"],
        ["7", "Stale JWT after profile update",     "Force window.location.reload() after save"],
        ["8", "Cron endpoint security",             "Bearer CRON_API_KEY header authentication"],
    ],
    0.4, 1.4, 12.5, 5.6
)
footer(s)
print("✅ Slide 16")

# ── SLIDE 17: Future Enhancements ───────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Future Enhancements", "What's Next for MediScript-E")
cols_data = [
    ("⚡  Short Term", [
        "Real-time notifications (WebSockets)",
        "Doctor availability calendar",
        "Patient-doctor rating system",
        "Push notifications (PWA)",
    ]),
    ("🚀  Medium Term", [
        "Video consultation (WebRTC)",
        "Prescription QR code for pharmacies",
        "Mobile app (React Native)",
        "Multi-language support (Bangla)",
    ]),
    ("🌐  Long Term", [
        "AI symptom checker",
        "Payment gateway for telemedicine",
        "HIPAA / GDPR compliance",
        "National health record integration",
    ]),
]
for i, (title, items) in enumerate(cols_data):
    x = 0.4 + i * 4.3
    rect(s, x, 1.4, 4.1, 5.5, fill=LIGHT_BG)
    rect(s, x, 1.4, 4.1, 0.55, fill=PRIMARY)
    tb(s, title, x+0.15, 1.48, 3.8, 0.38, size=14, bold=True, color=WHITE)
    bullets(s, items, x+0.15, 2.1, 3.8, 3.5, size=13)
footer(s)
print("✅ Slide 17")

# ── SLIDE 18: Conclusion ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Conclusion", "SDLC Complete — Results & Takeaway")
rect(s, 0.4, 1.4, 12.5, 0.55, fill=ACCENT)
tb(s, "SDLC Phases:   Requirements  →  Design  →  Implementation  →  Testing  →  Deployment   ✅",
   0.6, 1.5, 12.1, 0.38, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

achieved = [
    "Secure auth — Email Verification, 2FA, OAuth, RBAC",
    "Complete patient-doctor workflow — booking to prescription",
    "Automated medicine reminder system (cron + email)",
    "Secure cloud-based medical record storage",
    "AI-powered chatbot assistant (MediBot — Groq Llama 3.1)",
    "Deployed and live on Vercel — production ready",
]
tb(s, "What Was Achieved", 0.4, 2.1, 7.5, 0.4, size=15, bold=True, color=PRIMARY)
for i, item in enumerate(achieved):
    y = 2.6 + i * 0.55
    rect(s, 0.4, y, 0.4, 0.4, fill=GREEN)
    tb(s, "✓", 0.4, y+0.05, 0.4, 0.32, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, item, 0.9, y+0.05, 7.0, 0.35, size=13, color=DARK)

rect(s, 0.4, 6.0, 12.5, 1.0, fill=LIGHT_BG)
tb(s, '"Structured Software Engineering methodology turns ideas into maintainable,',
   0.6, 6.08, 12.1, 0.38, size=13, italic=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, 'scalable, and secure real-world applications."',
   0.6, 6.45, 12.1, 0.38, size=13, italic=True, color=PRIMARY, align=PP_ALIGN.CENTER)
footer(s)
print("✅ Slide 18")

# ── SLIDE 19: Q&A ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, fill=PRIMARY)
rect(s, 0, 0, 13.33, 7.5, fill=ACCENT)
rect(s, 0, 0, 13.33, 0.08, fill=WHITE)
rect(s, 0, 7.42, 13.33, 0.08, fill=WHITE)

tb(s, "Thank You", 0, 1.2, 13.33, 1.0, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "MediScript-E  —  Your Health, Digitally Managed.", 0, 2.3, 13.33, 0.55,
   size=18, italic=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

rect(s, 4.5, 3.1, 4.33, 0.05, fill=WHITE)

lines(s, [
    "Tushar Barua",
    "Computer Science & Engineering  ·  USTC, Chittagong",
], 0, 3.3, 13.33, 0.9, size=15, color=WHITE, align=PP_ALIGN.CENTER)

lines(s, [
    "📧  tusharcoder269@gmail.com",
    "🔗  github.com/tusharsno",
    "🌐  mediscript-e.vercel.app",
], 0, 4.35, 13.33, 1.2, size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

rect(s, 4.5, 5.75, 4.33, 0.05, fill=WHITE)
tb(s, "Open for Questions & Discussion", 0, 5.95, 13.33, 0.5,
   size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Video placeholder
rect(s, 9.5, 1.3, 3.5, 2.2, fill=RGBColor(0x0A, 0x3A, 0x50), line=WHITE)
tb(s, "🎬  PLAY VIDEO\nMediScript-E Demo", 9.5, 1.3, 3.5, 2.2,
   size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(s)
print("✅ Slide 19")

prs.save("MediScript-E-Presentation.pptx")
print("\n🎉 ALL 19 SLIDES COMPLETE — MediScript-E-Presentation.pptx ready!")
