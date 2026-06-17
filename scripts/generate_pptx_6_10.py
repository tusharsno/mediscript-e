from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY  = RGBColor(0x1A, 0x60, 0x80)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1E, 0x29, 0x3B)
GRAY     = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT   = RGBColor(0x0D, 0x4A, 0x63)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)

prs = Presentation("MediScript-E-Presentation.pptx")
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def lines(slide, items, x, y, w, h, size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return t

def bullets(slide, items, x, y, w, h, size=15, color=DARK):
    return lines(slide, [f"•  {i}" for i in items], x, y, w, h, size=size, color=color)

def img_placeholder(slide, x, y, w, h, label="ADD IMAGE"):
    b = rect(slide, x, y, w, h, fill=LIGHT_BG, line=PRIMARY)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"📸  {label}"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = PRIMARY
    return b

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.25, fill=PRIMARY)
    tb(slide, title, 0.4, 0.12, 12, 0.7, size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC", 0.3, 7.22, 12.7, 0.25,
       size=9, color=WHITE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    cw = Inches(w / cols)
    for i in range(cols): tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.color.rgb = DARK
    return tbl

# ── SLIDE 6: Use Cases ──────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Use Cases & Actors", "Phase 1 — Requirements Analysis")
rect(s, 0.4, 1.4, 5.5, 5.5, fill=LIGHT_BG)
tb(s, "Actors", 0.6, 1.5, 5.1, 0.4, size=15, bold=True, color=PRIMARY)
lines(s, [
    "👤  Patient",
    "👨⚕️  Doctor",
    "🛡️  Admin",
    "⚙️  Cron System (automated)",
    "🔗  OAuth Provider (Google / GitHub)",
], 0.6, 1.95, 5.1, 2.0, size=14)
tb(s, "Appointment Flow", 0.6, 4.1, 5.1, 0.4, size=14, bold=True, color=PRIMARY)
rect(s, 0.6, 4.55, 5.1, 0.6, fill=ACCENT)
tb(s, "PENDING  →  CONFIRMED  →  COMPLETED  |  CANCELLED",
   0.7, 4.65, 4.9, 0.4, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 6.3, 1.4, 6.6, 5.5, fill=LIGHT_BG)
tb(s, "Key Use Cases", 6.5, 1.5, 6.2, 0.4, size=15, bold=True, color=PRIMARY)
lines(s, [
    "UC1:  Register → Verify Email → Login",
    "UC2:  Patient books appointment",
    "UC3:  Doctor confirms & issues prescription",
    "UC4:  Patient downloads prescription PDF",
    "UC5:  Patient sets reminder → receives email",
    "UC6:  Admin manages users & monitors platform",
], 6.5, 1.95, 6.2, 3.5, size=13)
img_placeholder(s, 6.5, 5.5, 6.2, 1.2, "Use case diagram")
footer(s)
print("✅ Slide 6")

# ── SLIDE 7: System Architecture ────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "System Architecture", "Phase 2 — System Design")
layers = [
    (1.4,  PRIMARY, "PRESENTATION TIER",  "Next.js App Router  ·  React 19  ·  TypeScript"),
    (2.65, ACCENT,  "APPLICATION TIER",   "Next.js Serverless API Routes  ·  NextAuth.js  ·  Prisma ORM"),
    (3.9,  RGBColor(0x0A, 0x3A, 0x50), "DATA TIER", "PostgreSQL (Supabase)  ·  Supabase Storage"),
]
for y, color, title, sub in layers:
    rect(s, 0.5, y, 8.5, 1.0, fill=color)
    tb(s, title, 0.7, y+0.08, 8.0, 0.38, size=14, bold=True, color=WHITE)
    tb(s, sub,   0.7, y+0.52, 8.0, 0.38, size=12, color=LIGHT_BLUE)
    if y < 3.9:
        tb(s, "↓", 4.5, y+1.0, 0.5, 0.35, size=18, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

rect(s, 0.5, 5.15, 8.5, 0.05, fill=GRAY)
tb(s, "External Services", 0.5, 5.3, 8.5, 0.35, size=13, bold=True, color=PRIMARY)
lines(s, [
    "Groq API (Llama 3.1)  ·  Nodemailer (Gmail SMTP)  ·  Google & GitHub OAuth  ·  GitHub Actions (Cron)"
], 0.5, 5.65, 8.5, 0.4, size=12, color=GRAY)
img_placeholder(s, 9.3, 1.4, 3.7, 5.5, "Architecture diagram")
footer(s)
print("✅ Slide 7")

# ── SLIDE 8: Database Design ─────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Database Design", "Phase 2 — System Design")
rect(s, 0.4, 1.4, 7.8, 5.5, fill=LIGHT_BG)
tb(s, "Schema Overview  —  9 Models · 3 Roles · Cascade Delete", 0.6, 1.5, 7.4, 0.4, size=14, bold=True, color=PRIMARY)
lines(s, [
    "User",
    "    ├──  DoctorProfile   →  Appointment[ ]  ·  Prescription[ ]",
    "    └──  PatientProfile  →  Appointment[ ]  ·  Prescription[ ]",
    "                              MedicalVault[ ]  ·  MedicineReminder[ ]",
    "",
    "ContactMessage  (standalone)",
    "Testimonial     (standalone)",
], 0.6, 2.0, 7.4, 2.8, size=13, color=DARK)
tb(s, "Key Design Decisions", 0.6, 4.85, 7.4, 0.4, size=14, bold=True, color=PRIMARY)
bullets(s, [
    "Soft archive on prescriptions  (archivedByDoctor flag)",
    "OTP stored temporarily — cleared after verification",
    "Cascade delete — no orphan records on user deletion",
], 0.6, 5.3, 7.4, 1.3, size=13)
img_placeholder(s, 8.6, 1.4, 4.4, 5.5, "ER Diagram")
footer(s)
print("✅ Slide 8")

# ── SLIDE 9: Tech Stack ──────────────────────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Tech Stack", "Phase 3 — Implementation")
table(s,
    ["Layer", "Technology", "Purpose"],
    [
        ["Framework",    "Next.js 16 (App Router)",     "Full-stack web framework"],
        ["Language",     "TypeScript 5",                "Type safety across codebase"],
        ["Styling",      "Tailwind CSS 4 + Framer Motion", "UI & animations"],
        ["Auth",         "NextAuth.js 4",               "JWT sessions + OAuth"],
        ["ORM",          "Prisma 7",                    "Type-safe DB queries"],
        ["Database",     "PostgreSQL via Supabase",     "Managed relational DB"],
        ["Storage",      "Supabase Storage",            "Medical file uploads"],
        ["Email",        "Nodemailer (Gmail SMTP)",     "Verification, OTP, reminders"],
        ["AI",           "Groq SDK — Llama 3.1 8B",    "MediBot chatbot"],
        ["Deployment",   "Vercel",                      "Serverless hosting + CI/CD"],
    ],
    0.5, 1.4, 12.4, 5.6
)
footer(s)
print("✅ Slide 9")

# ── SLIDE 10: Authentication & Security ─────────────────
s = prs.slides.add_slide(BLANK)
header(s, "Authentication & Security", "Phase 3 — Implementation")
steps = [
    ("1", "Email Verification",  "crypto token · 24h expiry · account inactive until verified"),
    ("2", "Password Hashing",    "bcryptjs · 10 salt rounds"),
    ("3", "Two-Factor Auth (2FA)","email OTP · 6 digits · 10 min expiry · optional per user"),
    ("4", "OAuth Login",         "Google & GitHub · auto email-verified · profile auto-created"),
    ("5", "JWT Sessions",        "30-day expiry · role stored in token"),
    ("6", "RBAC",                "every API route checks role + ownership on PATCH/DELETE"),
]
for i, (num, title, desc) in enumerate(steps):
    y = 1.45 + i * 0.88
    rect(s, 0.4, y, 0.55, 0.6, fill=PRIMARY)
    tb(s, num, 0.4, y+0.08, 0.55, 0.45, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.05, y, 8.5, 0.6, fill=LIGHT_BG)
    tb(s, title, 1.15, y+0.05, 3.5, 0.3, size=14, bold=True, color=PRIMARY)
    tb(s, desc,  1.15, y+0.32, 8.2, 0.25, size=12, color=GRAY)
img_placeholder(s, 9.9, 1.45, 3.1, 2.5, "2FA screen screenshot")
img_placeholder(s, 9.9, 4.1,  3.1, 2.5, "Security architecture diagram")
footer(s)
print("✅ Slide 10")

prs.save("MediScript-E-Presentation.pptx")
print("✅ Slides 6–10 saved.")
