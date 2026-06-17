from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy
from lxml import etree

PRIMARY    = RGBColor(0x1A, 0x60, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
GRAY       = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG   = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT     = RGBColor(0x0D, 0x4A, 0x63)
LIGHT_BLUE = RGBColor(0xB0, 0xD4, 0xE3)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation("MediScript-E-Presentation.pptx")
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return t

def lines(slide, items, x, y, w, h, size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = True
    tf = t.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = item
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return t

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.25, fill=PRIMARY)
    tb(slide, title, 0.4, 0.12, 12, 0.7, size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.78, 12, 0.38, size=13, color=LIGHT_BLUE)

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill=PRIMARY)
    tb(slide, "MediScript-E  |  Tushar Barua  |  CSE, USTC", 0.3, 7.22, 12.7, 0.25,
       size=9, color=WHITE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    cw = Inches(w / cols)
    for i in range(cols): tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = hdr
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if ri % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(12); r.font.color.rgb = DARK
    return tbl

def move_slide(prs, old_index, new_index):
    """Move a slide from old_index to new_index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide = slides[old_index]
    xml_slides.remove(slide)
    xml_slides.insert(new_index, slide)

# ── BUILD OUTLINE SLIDE (will be moved to position 2) ──
outline = prs.slides.add_slide(BLANK)
header(outline, "Presentation Outline", "What We'll Cover Today")

outline_items = [
    ("1",  "Problem Statement",        "—"),
    ("2",  "Proposed Solution",        "—"),
    ("3",  "Requirements Analysis",    "Phase 1"),
    ("4",  "Use Cases & Actors",       "Phase 1"),
    ("5",  "System Architecture",      "Phase 2"),
    ("6",  "Database Design",          "Phase 2"),
    ("7",  "Tech Stack",               "Phase 3"),
    ("8",  "Authentication & Security","Phase 3"),
    ("9",  "Core Features",            "Phase 3"),
    ("10", "Testing",                  "Phase 4"),
    ("11", "Results & Evaluation",     "Phase 5"),
    ("12", "Deployment",               "Phase 5"),
    ("13", "Challenges & Future Work", "—"),
    ("14", "Conclusion",               "—"),
]

# Two columns
left  = outline_items[:7]
right = outline_items[7:]

for i, (num, title, phase) in enumerate(left):
    y = 1.45 + i * 0.77
    rect(outline, 0.4, y, 0.5, 0.58, fill=PRIMARY)
    tb(outline, num, 0.4, y+0.1, 0.5, 0.38, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(outline, 0.95, y, 5.5, 0.58, fill=LIGHT_BG if i % 2 == 0 else WHITE)
    tb(outline, title, 1.05, y+0.1, 3.8, 0.38, size=13, color=DARK)
    tb(outline, phase, 4.85, y+0.1, 1.5, 0.38, size=11, color=PRIMARY, bold=True, align=PP_ALIGN.RIGHT)

for i, (num, title, phase) in enumerate(right):
    y = 1.45 + i * 0.77
    rect(outline, 6.9, y, 0.5, 0.58, fill=PRIMARY)
    tb(outline, num, 6.9, y+0.1, 0.5, 0.38, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(outline, 7.45, y, 5.5, 0.58, fill=LIGHT_BG if i % 2 == 0 else WHITE)
    tb(outline, title, 7.55, y+0.1, 3.8, 0.38, size=13, color=DARK)
    tb(outline, phase, 11.35, y+0.1, 1.5, 0.38, size=11, color=PRIMARY, bold=True, align=PP_ALIGN.RIGHT)

footer(outline)
# Move outline slide from last position (index 19) to index 2 (after Abstract)
move_slide(prs, len(prs.slides)-1, 2)
print("✅ Outline slide injected at position 3")

# ── BUILD RESULTS & EVALUATION SLIDE (will be moved to position 15) ──
results = prs.slides.add_slide(BLANK)
header(results, "Results & Evaluation", "Phase 5 — What Was Delivered")

delivered = [
    "Email Verification (24h token)",
    "Two-Factor Authentication (2FA)",
    "Google & GitHub OAuth",
    "Appointment Booking & Management",
    "Digital Prescriptions + PDF Download",
    "Automated Medicine Reminders (email)",
    "Medical Vault (cloud storage)",
    "AI Chatbot — MediBot (Groq Llama 3.1)",
    "Global Search (role-aware, Ctrl+K)",
    "Admin Dashboard & User Management",
    "Responsive UI — mobile + desktop",
    "Production Deployment — Live on Vercel",
]

left_items  = delivered[:6]
right_items = delivered[6:]

for i, item in enumerate(left_items):
    y = 1.45 + i * 0.88
    rect(results, 0.4, y, 0.45, 0.65, fill=GREEN)
    tb(results, "✓", 0.4, y+0.12, 0.45, 0.4, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(results, 0.9, y, 5.6, 0.65, fill=LIGHT_BG if i % 2 == 0 else WHITE)
    tb(results, item, 1.0, y+0.15, 5.4, 0.38, size=13, color=DARK)

for i, item in enumerate(right_items):
    y = 1.45 + i * 0.88
    rect(results, 6.9, y, 0.45, 0.65, fill=GREEN)
    tb(results, "✓", 6.9, y+0.12, 0.45, 0.4, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(results, 7.4, y, 5.6, 0.65, fill=LIGHT_BG if i % 2 == 0 else WHITE)
    tb(results, item, 7.5, y+0.15, 5.4, 0.38, size=13, color=DARK)

rect(results, 0.4, 6.85, 12.5, 0.22, fill=ACCENT)
tb(results, "All SDLC phases completed  ·  All requirements met  ·  Live on Vercel",
   0.4, 6.87, 12.5, 0.2, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

footer(results)
# Current last index = 20, move to index 15 (after Testing which is at index 14)
move_slide(prs, len(prs.slides)-1, 15)
print("✅ Results & Evaluation slide injected at position 16")

prs.save("MediScript-E-Presentation.pptx")
print(f"\n✅ Total slides: {len(prs.slides)}")
print("🎉 MediScript-E-Presentation.pptx updated — 21 slides ready!")
