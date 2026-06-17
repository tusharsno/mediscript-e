from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import os

SS = "public/MediScript-E_SS/"

prs = Presentation("MediScript-E-Presentation.pptx")

def remove_placeholders(slide):
    """Remove all existing 📸 placeholder shapes from a slide."""
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "📸" in run.text or "ADD IMAGE" in run.text:
                        to_remove.append(shape)
                        break
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

def add_image(slide, path, x, y, w, h):
    """Add image to slide at given position."""
    if not os.path.exists(path):
        print(f"  ⚠️  File not found: {path}")
        return
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def add_image_fit(slide, path, x, y, w, h):
    """Add image maintaining aspect ratio within given box."""
    if not os.path.exists(path):
        print(f"  ⚠️  File not found: {path}")
        return
    from PIL import Image as PILImage
    img = PILImage.open(path)
    iw, ih = img.size
    ratio = iw / ih
    box_ratio = w / h
    if ratio > box_ratio:
        new_w = w
        new_h = w / ratio
        y_offset = y + (h - new_h) / 2
        x_offset = x
    else:
        new_h = h
        new_w = h * ratio
        x_offset = x + (w - new_w) / 2
        y_offset = y
    slide.shapes.add_picture(path, Inches(x_offset), Inches(y_offset), Inches(new_w), Inches(new_h))

# ── SLIDE 1: Title ──────────────────────────────────────
# Right side image placeholder (x=7.3, y=1.0, w=5.7, h=5.5)
s = prs.slides[0]
remove_placeholders(s)
add_image_fit(s, SS + "Landing_page_1.png", 7.3, 1.0, 5.7, 5.5)
print("✅ Slide 1 — Landing_page_1.png")

# ── SLIDE 2: Abstract ───────────────────────────────────
# Right side (x=8.8, y=1.45, w=4.1, h=4.5)
s = prs.slides[1]
remove_placeholders(s)
add_image_fit(s, SS + "Patient-Dashboard-overview-1.png", 8.8, 1.45, 4.1, 4.5)
print("✅ Slide 2 — Patient-Dashboard-overview-1.png")

# ── SLIDE 7: Use Cases ──────────────────────────────────
# Bottom placeholder (x=6.5, y=5.5, w=6.2, h=1.2) — use landing page section
s = prs.slides[6]
remove_placeholders(s)
add_image_fit(s, SS + "Landing-page-2.png", 6.5, 5.3, 6.2, 1.55)
print("✅ Slide 7 — Landing-page-2.png (use case area)")

# ── SLIDE 8: Architecture ───────────────────────────────
# Right side (x=9.3, y=1.4, w=3.7, h=5.5)
s = prs.slides[7]
remove_placeholders(s)
add_image_fit(s, SS + "Landing-page-3.png", 9.3, 1.4, 3.7, 5.5)
print("✅ Slide 8 — Landing-page-3.png (architecture area)")

# ── SLIDE 9: Database ───────────────────────────────────
# Right side (x=8.6, y=1.4, w=4.4, h=5.5)
s = prs.slides[8]
remove_placeholders(s)
add_image_fit(s, SS + "Landing-page-4.png", 8.6, 1.4, 4.4, 5.5)
print("✅ Slide 9 — Landing-page-4.png (DB area)")

# ── SLIDE 11: Auth & Security ───────────────────────────
# Two placeholders: 2FA screen (x=9.9, y=1.45, w=3.1, h=2.5)
#                   Security diagram (x=9.9, y=4.1, w=3.1, h=2.5)
s = prs.slides[10]
remove_placeholders(s)
add_image_fit(s, SS + "2FA-OTP-screen.png",            9.9, 1.45, 3.1, 2.5)
add_image_fit(s, SS + "Landing-page-05.png",           9.9, 4.1,  3.1, 2.5)
print("✅ Slide 11 — 2FA-OTP-screen.png + Landing-page-05.png")

# ── SLIDE 12: Patient Features ──────────────────────────
# 4 feature cards each with a bottom image placeholder
# Card positions: (0.4,1.4), (6.5,1.4), (0.4,4.05), (6.5,4.05)
# Each placeholder: x+0.2, y+1.3, w=5.7, h=0.85
s = prs.slides[11]
remove_placeholders(s)
imgs_patient = [
    (SS + "Patient-Dashboard-overview-1.png", 0.6,  2.75, 5.7, 0.85),
    (SS + "Landing-page-6.png",               6.7,  2.75, 5.7, 0.85),
    (SS + "Landing-page-7.png",               0.6,  5.4,  5.7, 0.85),
    (SS + "Landing-page-8.png",               6.7,  5.4,  5.7, 0.85),
]
for path, x, y, w, h in imgs_patient:
    add_image_fit(s, path, x, y, w, h)
print("✅ Slide 12 — Patient dashboard + 3 landing page sections")

# ── SLIDE 13: Doctor & Admin ────────────────────────────
# Left: Doctor (x=0.6, y=4.15, w=5.7, h=2.5)
# Right: Admin (x=7.1, y=3.65, w=5.7, h=3.0)
s = prs.slides[12]
remove_placeholders(s)
add_image_fit(s, SS + "Doctor-Dashboard-overview-1.png", 0.6, 4.15, 5.7, 2.5)
add_image_fit(s, SS + "Admin-Dashboard-overview-1.png",  7.1, 3.65, 5.7, 3.0)
print("✅ Slide 13 — Doctor-Dashboard + Admin-Dashboard")

# ── SLIDE 14: MediBot & Search ──────────────────────────
# Left bottom (x=0.6, y=4.45, w=5.7, h=2.2)
s = prs.slides[13]
remove_placeholders(s)
add_image_fit(s, SS + "MediBot.png", 0.6, 4.45, 5.7, 2.2)
print("✅ Slide 14 — MediBot.png")

# ── SLIDE 17: Deployment ────────────────────────────────
# Right side (x=8.3, y=1.45, w=4.7, h=5.5)
s = prs.slides[16]
remove_placeholders(s)
add_image_fit(s, SS + "Vercel-deployment-dashboard.png", 8.3, 1.45, 4.7, 5.5)
print("✅ Slide 17 — Vercel-deployment-dashboard.png")

prs.save("MediScript-E-Presentation.pptx")
print("\n🎉 All screenshots injected! MediScript-E-Presentation.pptx updated.")
